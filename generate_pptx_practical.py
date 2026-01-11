import os
import time
import requests
from qdrant_client import QdrantClient
from pptx import Presentation

# ----------------- ENV -----------------
PROJECT_ID = os.environ["PROJECT_ID"]              # odontology-rag-slides
EMB_LOCATION = os.environ.get("LOCATION", "europe-west1")  # embeddings region
GEN_LOCATION = os.environ.get("GEN_LOCATION", "global")    # generation location

QDRANT_URL = os.environ["QDRANT_URL"]
QDRANT_API_KEY = os.environ["QDRANT_API_KEY"]
COL = os.environ["QDRANT_COLLECTION"]

# Models
EMB_MODEL = "gemini-embedding-001"
GEN_MODEL = os.environ.get("GEN_MODEL_ID", "gemini-2.5-flash")  # practical + fast

# ----------------- Google token -----------------
def gcloud_token() -> str:
    tok = os.popen("gcloud auth print-access-token").read().strip()
    if not tok:
        raise RuntimeError("No gcloud access token. Run: gcloud auth login")
    return tok

# ----------------- Vertex Embedding (regional predict endpoint works) -----------------
def embed(text: str):
    url = (
        f"https://{EMB_LOCATION}-aiplatform.googleapis.com/v1/projects/"
        f"{PROJECT_ID}/locations/{EMB_LOCATION}/publishers/google/models/{EMB_MODEL}:predict"
    )
    r = requests.post(
        url,
        headers={"Authorization": f"Bearer {gcloud_token()}", "Content-Type": "application/json"},
        json={"instances": [{"content": text}]},
        timeout=60,
    )
    r.raise_for_status()
    return r.json()["predictions"][0]["embeddings"]["values"]

# ----------------- Vertex Generation (USE GLOBAL ENDPOINT) -----------------
def generate_bullets(prompt: str) -> list[str]:
    api_endpoint = "https://aiplatform.googleapis.com"
    url = (
        f"{api_endpoint}/v1/projects/{PROJECT_ID}/locations/{GEN_LOCATION}/publishers/"
        f"google/models/{GEN_MODEL}:generateContent"
    )
    r = requests.post(
        url,
        headers={"Authorization": f"Bearer {gcloud_token()}", "Content-Type": "application/json"},
        json={
            "contents": [{"role": "user", "parts": [{"text": prompt}]}],
            "generationConfig": {"temperature": 0.25, "maxOutputTokens": 320},
        },
        timeout=120,
    )
    r.raise_for_status()
    text = r.json()["candidates"][0]["content"]["parts"][0]["text"].strip()

    # Normalize to bullet list
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    bullets = []
    for ln in lines:
        ln = ln.lstrip("-•").strip()
        if ln:
            bullets.append(ln)
    return bullets[:4] if bullets else ["(No content generated)"]

# ----------------- Qdrant retrieval -----------------
def retrieve(query: str, limit: int = 24):
    q = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY)
    vec = embed(query)
    res = q.query_points(collection_name=COL, query=vec, limit=limit)
    return res.points

def format_citations(hits, max_items=6) -> str:
    cites = []
    for h in hits[:max_items]:
        p = h.payload or {}
        cites.append(f"- {p.get('source_object_path')} (pp. {p.get('page_start')}-{p.get('page_end')})")
    return "Sources used:\n" + "\n".join(cites)

# ----------------- PPT helpers -----------------
def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: str):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for b in bullets:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = b
    slide.notes_slide.notes_text_frame.text = notes

# ----------------- Main -----------------
def main():
    talk_title = "Dental Adhesive Systems — Practical Clinical Update"
    audience = "practicing dentists / restorative colleagues"
    base_query = "dental adhesive systems etch-and-rinse self-etch universal clinical performance meta-analysis technique sensitivity"

    outline = [
        ("Why this matters in daily practice", "clinical relevance, failure modes, marginal integrity, post-op sensitivity"),
        ("System taxonomy (what to choose)", "etch-and-rinse vs self-etch vs universal; when to use which"),
        ("Etch-and-rinse: best indications", "enamel bonding, selective etch, control of moisture, pitfalls"),
        ("Self-etch: when it shines / fails", "dentin performance, mild vs strong, enamel limitations"),
        ("Universal adhesives: practical protocol", "selective enamel etch, active application, solvent evaporation"),
        ("Evidence snapshot: clinical outcomes", "retention, marginal discoloration, post-op sensitivity, longevity"),
        ("Technique sensitivity: common errors", "contamination, under-etch, pooling, inadequate air thinning, curing"),
        ("Enamel vs dentin: chairside tips", "selective etch, smear layer, dentin wetness, scrubbing time"),
        ("Longevity & repairs", "when to repair, surface prep, re-etch, re-adhere, maintenance"),
        ("Take-home checklist", "simple decision rules and protocol checklist"),
    ]

    prs = Presentation()
    add_title_slide(
        prs,
        talk_title,
        f"Practical, evidence-informed slides for {audience}\n(Generated from your indexed PDFs via Qdrant RAG)"
    )

    for (slide_title, slide_focus) in outline:
        # retrieve slide-specific evidence
        query = f"{base_query}. Focus: {slide_focus}"
        hits = retrieve(query, limit=18)

        # build short context from top hits
        ctx_parts = []
        for h in hits[:10]:
            p = h.payload or {}
            t = (p.get("text") or "").replace("\n", " ").strip()
            if t:
                ctx_parts.append(t[:400])
        context = "\n".join(ctx_parts)[:6000]

        prompt = f"""
You are a dental expert. Create content for ONE PowerPoint slide.

Audience: {audience}
Language: English
Tone: practical, chairside, concise (no theory lecture)
Slide title: {slide_title}

Task:
- Write 3–4 bullet points (max 14 words each).
- Focus on actionable clinical recommendations and pitfalls.
- Avoid citations in bullets (citations go in speaker notes).

Evidence context (from the user's documents):
{context}
""".strip()

        bullets = generate_bullets(prompt)
        notes = format_citations(hits, max_items=7)
        add_bullet_slide(prs, slide_title, bullets, notes)

        time.sleep(0.2)

    out = "Dental_Adhesive_Systems_Practical_10slides.pptx"
    prs.save(out)
    print(f"✅ Saved: {out}")

if __name__ == "__main__":
    main()

