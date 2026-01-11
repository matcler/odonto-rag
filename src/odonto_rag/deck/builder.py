from __future__ import annotations

import argparse
import json
import math
import os
import re
import sys
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER


# ============================================================
# CONFIG
# ============================================================

LABEL_ORDER = ["DO", "AVOID", "CHECK", "RESCUE", "TIP"]

LABEL_COLORS_FALLBACK = {
    "DO": RGBColor(0, 128, 0),
    "AVOID": RGBColor(192, 0, 0),
    "CHECK": RGBColor(0, 102, 204),
    "RESCUE": RGBColor(128, 0, 128),
    "TIP": RGBColor(102, 102, 102),
}
FALLBACK_FONT_NAME = "Calibri"

TITLE_FONT_SIZE = 36
BULLET_FONT_SIZES = [22, 20, 18, 16]
PARA_SPACE_AFTER_PT = 8

TRUNCATE_CHAR_LIMIT = 220
TRUNCATE_SUFFIX = " …"

EMU_PER_INCH = 914400

# Preview parsing (more permissive)
# Accept:
#   Slide 1 : Title
#   Slide 1: Title
#   Slide 1 - Title
#   1 - Title
_PREVIEW_SLIDE_RE = re.compile(r"^\s*(?:Slide\s*)?(\d+)\s*[:\-]\s*(.+?)\s*$", re.IGNORECASE)

# Accept bullets:
#   - Do: text
#   • Do: text
#   Do: text
_PREVIEW_BULLET_RE = re.compile(r"^\s*(?:[-•]\s*)?([A-Za-z']+)\s*:\s*(.+?)\s*$")

# JSON deck bullets as "Label: text"
BULLET_LINE_RE = re.compile(r"^\s*([A-Za-z']+)\s*:\s*(.+?)\s*$")


class SlideSpecError(Exception):
    pass


# ============================================================
# Slide formats (ONLY when no template)
# ============================================================

_FORMATS_INCH = {
    "widescreen": (13.333, 7.5),
    "standard": (10.0, 7.5),
    "a4landscape": (11.69, 8.27),
    "a4portrait": (8.27, 11.69),
}

def _parse_format(fmt: str) -> Tuple[float, float]:
    fmt = (fmt or "").strip().lower()
    if not fmt:
        return _FORMATS_INCH["widescreen"]
    if fmt in _FORMATS_INCH:
        return _FORMATS_INCH[fmt]
    if fmt.startswith("custom:"):
        rest = fmt.split("custom:", 1)[1]
        m = re.match(r"^\s*([0-9]+(\.[0-9]+)?)x([0-9]+(\.[0-9]+)?)\s*$", rest)
        if not m:
            raise SlideSpecError("Invalid custom format. Use custom:WIDTHxHEIGHT in inches, e.g. custom:13.333x7.5")
        return float(m.group(1)), float(m.group(3))
    raise SlideSpecError(f"Unknown format '{fmt}'")


# ============================================================
# Layout computed from slide size (percent-based fallback)
# ============================================================

@dataclass(frozen=True)
class Layout:
    title_left: float
    title_top: float
    title_width: float
    title_height: float
    body_left: float
    body_top: float
    body_width: float
    body_height: float

def _emu_to_inches(x_emu: int) -> float:
    return float(x_emu) / EMU_PER_INCH

def _compute_layout(prs: Presentation) -> Layout:
    w = _emu_to_inches(prs.slide_width)
    h = _emu_to_inches(prs.slide_height)
    return Layout(
        0.05 * w, 0.05 * h, 0.90 * w, 0.12 * h,
        0.06 * w, 0.20 * h, 0.88 * w, 0.72 * h
    )


# ============================================================
# Templates: list + resolve
# ============================================================

def list_templates(templates_dir: str) -> List[str]:
    d = templates_dir or "templates"
    if not os.path.isdir(d):
        return []
    out = []
    for name in os.listdir(d):
        if name.lower().endswith(".pptx") and not name.startswith("~$"):
            out.append(name)
    out.sort(key=lambda s: s.lower())
    return out

def list_templates_payload(templates_dir: str) -> Dict[str, Any]:
    d = templates_dir or "templates"
    names = list_templates(d)
    return {"templates_dir": d, "templates": [{"name": n, "path": os.path.join(d, n)} for n in names]}

def resolve_template_path(template: Optional[str], template_name: Optional[str], templates_dir: str) -> Optional[str]:
    if template:
        return template
    if template_name:
        return os.path.join(templates_dir or "templates", template_name)
    return None


# ============================================================
# Helpers
# ============================================================

def _label_norm(label: str) -> str:
    l = (label or "").strip().upper()
    if l in ("DO",): return "DO"
    if l in ("AVOID", "DONT", "DON'T"): return "AVOID"
    if l in ("CHECK",): return "CHECK"
    if l in ("RESCUE", "FIX"): return "RESCUE"
    if l in ("TIP", "PEARL"): return "TIP"
    return l


# ============================================================
# Input parsing (JSON OR preview OR stdin)
# ============================================================

def _read_text_input(inp: str) -> str:
    if inp in ("/dev/stdin", "-"):
        return sys.stdin.read()
    with open(inp, "r", encoding="utf-8") as f:
        return f.read()

def _try_parse_json(text: str) -> Optional[dict]:
    try:
        obj = json.loads(text)
        return obj if isinstance(obj, dict) else None
    except Exception:
        return None

def parse_slide_writer_preview(text: str) -> dict:
    slides: List[Dict[str, Any]] = []
    cur: Optional[Dict[str, Any]] = None

    for raw_line in text.splitlines():
        line = raw_line.rstrip("\n")
        m = _PREVIEW_SLIDE_RE.match(line)
        if m:
            if cur:
                slides.append(cur)
            cur = {"title": m.group(2).strip(), "bullets": [], "notes": ""}
            continue

        if cur:
            bm = _PREVIEW_BULLET_RE.match(line)
            if bm:
                cur["bullets"].append({"label": _label_norm(bm.group(1)), "text": bm.group(2).strip()})

    if cur:
        slides.append(cur)

    # basic validation: must have at least one slide with a title
    if not slides or not any(s.get("title") for s in slides):
        raise SlideSpecError("Input is neither valid JSON deck nor slide preview (no slide headings found).")

    return {"deck_title": "", "language": "", "slides": slides}

def from_slide_writer_deck(deck: dict) -> dict:
    slides_out = []
    slides_in = deck.get("slides", [])
    if not isinstance(slides_in, list) or not slides_in:
        raise SlideSpecError("slide_writer deck has no slides.")

    for s in slides_in:
        title = str(s.get("title", "")).strip()
        if not title:
            continue

        parsed = []
        for line in s.get("bullets", []):
            m = BULLET_LINE_RE.match(str(line))
            if not m:
                continue
            parsed.append({"label": _label_norm(m.group(1)), "text": m.group(2).strip()})

        slides_out.append({"title": title, "bullets": parsed, "notes": str(s.get("notes", "") or "").strip()})

    return {"deck_title": deck.get("topic", ""), "language": "", "slides": slides_out}

def normalize_slidespec(spec: dict) -> dict:
    out = {"deck_title": spec.get("deck_title", ""), "language": "", "slides": []}
    for s in spec["slides"]:
        mapped: Dict[str, str] = {}
        for b in s.get("bullets", []):
            label = _label_norm(b.get("label"))
            text = (b.get("text", "") or "").strip()
            if not label:
                continue
            mapped[label] = (mapped.get(label, "") + " • " + text).strip(" •")
        bullets = [{"label": k, "text": mapped.get(k, "")} for k in LABEL_ORDER]
        out["slides"].append({"title": s["title"], "bullets": bullets, "notes": s.get("notes", "")})
    return out

def load_input_auto(inp: Union[str, Dict[str, Any]]) -> Tuple[dict, str]:
    if isinstance(inp, dict):
        return inp, "dict"
    text = _read_text_input(inp)
    if not text.strip():
        return {"deck_title": "", "language": "", "slides": []}, "empty"
    obj = _try_parse_json(text)
    if obj is not None:
        return obj, "json"
    return parse_slide_writer_preview(text), "preview"


# ============================================================
# Overflow logic
# ============================================================

def _estimate_lines(text: str, font_pt: int, width_in: float) -> int:
    if not text:
        return 1
    k = 0.55
    chars = int((width_in * 72 / font_pt) * k)
    return max(1, math.ceil(len(text) / max(chars, 12)))

def _fits(bullets, font_pt, body_width_in: float, body_height_in: float):
    max_h = body_height_in * 72
    used = 0
    for b in bullets:
        lines = _estimate_lines(f"{b['label']}: {b['text']}", font_pt, body_width_in)
        used += int(1.15 * font_pt) * lines + PARA_SPACE_AFTER_PT
    return used <= max_h

def _truncate(text: str):
    text = text or ""
    if len(text) <= TRUNCATE_CHAR_LIMIT:
        return text, False
    return text[: max(0, TRUNCATE_CHAR_LIMIT - len(TRUNCATE_SUFFIX))] + TRUNCATE_SUFFIX, True


# ============================================================
# Template slide removal
# ============================================================

def _delete_all_slides(prs: Presentation) -> None:
    sldIdLst = prs.slides._sldIdLst  # pylint: disable=protected-access
    if sldIdLst is None:
        return
    for sldId in list(sldIdLst):
        rId = sldId.rId
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


# ============================================================
# Find a good layout with placeholders
# ============================================================

def _find_title_body_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        has_title = False
        has_body = False
        for ph in layout.placeholders:
            pht = ph.placeholder_format.type
            if pht == PP_PLACEHOLDER.TITLE:
                has_title = True
            if pht in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                has_body = True
        if has_title and has_body:
            return layout
    if len(prs.slide_layouts) > 1:
        return prs.slide_layouts[1]
    return prs.slide_layouts[0]

def _get_placeholder(slide, phtype):
    for shape in slide.placeholders:
        if shape.placeholder_format.type == phtype:
            return shape
    return None


# ============================================================
# Rendering
# ============================================================

def _render_with_placeholders(
    slide,
    title: str,
    bullets: List[Dict[str, str]],
    chosen_font_pt: int,
    use_template_fonts: bool,
    use_template_colors: bool,
):
    title_ph = _get_placeholder(slide, PP_PLACEHOLDER.TITLE)
    body_ph = _get_placeholder(slide, PP_PLACEHOLDER.BODY) or _get_placeholder(slide, PP_PLACEHOLDER.OBJECT)
    if title_ph is None or body_ph is None:
        return False

    title_ph.text = title
    if not use_template_fonts:
        try:
            p = title_ph.text_frame.paragraphs[0]
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(TITLE_FONT_SIZE)
            r.font.bold = True
            r.font.name = FALLBACK_FONT_NAME
        except Exception:
            pass

    tf = body_ph.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = 0
        para.space_after = Pt(PARA_SPACE_AFTER_PT)

        rl = para.add_run()
        rl.text = f"{b['label']}: "
        rl.font.bold = True
        rl.font.size = Pt(chosen_font_pt)
        if not use_template_fonts:
            rl.font.name = FALLBACK_FONT_NAME
        if not use_template_colors:
            rl.font.color.rgb = LABEL_COLORS_FALLBACK[b["label"]]

        rt = para.add_run()
        rt.text = b["text"]
        rt.font.size = Pt(chosen_font_pt)
        if not use_template_fonts:
            rt.font.name = FALLBACK_FONT_NAME

    return True


def _render_fallback_textboxes(
    slide,
    layout: Layout,
    title: str,
    bullets: List[Dict[str, str]],
    chosen_font_pt: int,
    use_template_fonts: bool,
    use_template_colors: bool,
):
    tbox = slide.shapes.add_textbox(
        Inches(layout.title_left), Inches(layout.title_top),
        Inches(layout.title_width), Inches(layout.title_height)
    )
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.level = 0
    r = p.add_run()
    r.text = title
    if not use_template_fonts:
        r.font.name = FALLBACK_FONT_NAME
        r.font.size = Pt(TITLE_FONT_SIZE)
        r.font.bold = True

    box = slide.shapes.add_textbox(
        Inches(layout.body_left), Inches(layout.body_top),
        Inches(layout.body_width), Inches(layout.body_height)
    )
    btf = box.text_frame
    btf.clear()
    btf.word_wrap = True

    for i, b in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.level = 0
        para.space_after = Pt(PARA_SPACE_AFTER_PT)

        rl = para.add_run()
        rl.text = f"{b['label']}: "
        rl.font.bold = True
        rl.font.size = Pt(chosen_font_pt)
        if not use_template_fonts:
            rl.font.name = FALLBACK_FONT_NAME
        if not use_template_colors:
            rl.font.color.rgb = LABEL_COLORS_FALLBACK[b["label"]]

        rt = para.add_run()
        rt.text = b["text"]
        rt.font.size = Pt(chosen_font_pt)
        if not use_template_fonts:
            rt.font.name = FALLBACK_FONT_NAME


def _render_slide(
    prs: Presentation,
    slide_data: dict,
    layout: Layout,
    use_template_fonts: bool,
    use_template_colors: bool,
    base_layout,
):
    slide = prs.slides.add_slide(base_layout)

    bullets = slide_data["bullets"]
    chosen = next(
        (s for s in BULLET_FONT_SIZES if _fits(bullets, s, layout.body_width, layout.body_height)),
        BULLET_FONT_SIZES[-1],
    )

    trunc_map, full_notes = {}, []
    if not _fits(bullets, chosen, layout.body_width, layout.body_height):
        for b in bullets:
            t, cut = _truncate(b["text"])
            if cut:
                trunc_map[b["label"]] = t
                full_notes.append(f"{b['label']}: {b['text']}")

    bullets2 = [{"label": b["label"], "text": trunc_map.get(b["label"], b["text"])} for b in bullets]

    ok = _render_with_placeholders(
        slide, slide_data["title"], bullets2, chosen,
        use_template_fonts=use_template_fonts,
        use_template_colors=use_template_colors,
    )
    if not ok:
        _render_fallback_textboxes(
            slide, layout, slide_data["title"], bullets2, chosen,
            use_template_fonts=use_template_fonts,
            use_template_colors=use_template_colors,
        )

    if slide_data.get("notes") or full_notes:
        nt = slide.notes_slide.notes_text_frame
        nt.text = (slide_data.get("notes") or "").strip()
        if full_notes:
            if nt.text:
                nt.text += "\n\n"
            nt.text += "[TRUNCATED ON SLIDE]\n" + "\n".join(full_notes)


# ============================================================
# PUBLIC API
# ============================================================

def build_pptx(
    inp: Union[str, Dict[str, Any]],
    out_path: str,
    template: Optional[str] = None,
    template_name: Optional[str] = None,
    templates_dir: str = "templates",
    fmt: str = "widescreen",
):
    raw, kind = load_input_auto(inp)

    if kind == "empty":
        raise SlideSpecError("Input is empty (stdin or file had no content).")

    spec = from_slide_writer_deck(raw) if isinstance(raw, dict) and "topic" in raw else raw
    spec = normalize_slidespec(spec)

    slides_in = spec.get("slides", [])
    if not isinstance(slides_in, list) or not slides_in:
        raise SlideSpecError("No slides to render after parsing/normalization.")

    template_path = resolve_template_path(template, template_name, templates_dir)
    if template_path and not os.path.exists(template_path):
        raise SlideSpecError(f"Template not found: {template_path}")

    prs = Presentation(template_path) if template_path else Presentation()

    use_template_fonts = bool(template_path)
    use_template_colors = bool(template_path)

    if not template_path:
        w, h = _parse_format(fmt)
        prs.slide_width = Inches(w)
        prs.slide_height = Inches(h)

    base_layout = _find_title_body_layout(prs)

    if template_path and len(prs.slides) > 0:
        _delete_all_slides(prs)

    layout = _compute_layout(prs)

    for s in slides_in:
        _render_slide(prs, s, layout, use_template_fonts, use_template_colors, base_layout)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)

def inspect_input(inp: str, max_lines: int = 80) -> int:
    text = _read_text_input(inp)
    if not text.strip():
        print("[inspect] input is EMPTY")
        return 2

    obj = _try_parse_json(text)
    if obj is not None:
        keys = ", ".join(sorted(list(obj.keys()))[:20])
        print(f"[inspect] detected JSON dict. keys: {keys}")
        return 0

    # try preview detection: any slide headings?
    slide_lines = [ln for ln in text.splitlines() if _PREVIEW_SLIDE_RE.match(ln)]
    bullet_lines = [ln for ln in text.splitlines() if _PREVIEW_BULLET_RE.match(ln)]
    if slide_lines:
        print(f"[inspect] detected PREVIEW. slide_headings={len(slide_lines)} bullet_like={len(bullet_lines)}")
        return 0

    print("[inspect] UNKNOWN FORMAT. showing first lines:")
    for i, ln in enumerate(text.splitlines()[:max_lines], start=1):
        print(f"{i:02d}: {ln}")
    return 1


# ============================================================
# CLI
# ============================================================

def main():
    ap = argparse.ArgumentParser(description="Chairside PPTX builder (JSON or slide_writer preview).")
    sub = ap.add_subparsers(dest="cmd", required=True)

    lt = sub.add_parser("list-templates", help="List available PPTX templates in a directory.")
    lt.add_argument("--templates-dir", default="templates")
    lt.add_argument("--json", action="store_true", help="Print JSON payload instead of plain lines.")

    ins = sub.add_parser("inspect-in", help="Inspect input format (JSON vs preview) and print a safe snippet.")
    ins.add_argument("--in", dest="inp", required=True, help="Path or /dev/stdin or -")

    b = sub.add_parser("build", help="Build PPTX from JSON deck or slide_writer preview.")
    b.add_argument("--in", dest="inp", required=True, help="Path or /dev/stdin or -")
    b.add_argument("--out", required=True)
    b.add_argument("--templates-dir", default="templates")
    b.add_argument("--template", default=None, help="Explicit template path (.pptx).")
    b.add_argument("--template-name", default=None, help="Template filename inside --templates-dir.")
    b.add_argument("--format", default="widescreen", help="Used ONLY when no template is provided.")

    args = ap.parse_args()

    if args.cmd == "list-templates":
        if args.json:
            print(json.dumps(list_templates_payload(args.templates_dir), ensure_ascii=False, indent=2))
        else:
            for name in list_templates(args.templates_dir):
                print(name)
        return

    if args.cmd == "inspect-in":
        code = inspect_input(args.inp)
        raise SystemExit(code)

    build_pptx(
        args.inp,
        args.out,
        template=args.template,
        template_name=args.template_name,
        templates_dir=args.templates_dir,
        fmt=args.format,
    )

if __name__ == "__main__":
    main()
