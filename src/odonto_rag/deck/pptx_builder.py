from __future__ import annotations

import hashlib
import json
import os
from pathlib import Path
import re
import subprocess
from typing import Any, Dict, Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

_KEYNOTE_SAFE = str(os.environ.get("PPTX_KEYNOTE_SAFE") or "").strip().lower() in {"1", "true", "yes", "on"}
_KEYNOTE_KEEP_TABLES = str(os.environ.get("PPTX_KEYNOTE_KEEP_TABLES") or "").strip().lower() in {
    "1",
    "true",
    "yes",
    "on",
}
_ELLIPSIS = "..." if _KEYNOTE_SAFE else "…"

FONT_NAME = "Arial" if _KEYNOTE_SAFE else "Calibri"
TITLE_FONT_SIZE_PT = 32
BODY_FONT_SIZE_PT = 20
FOOTER_FONT_SIZE_PT = 8
MAX_BULLETS_PER_SLIDE = 6
MAX_BULLET_CHARS = 170
MAX_TITLE_CHARS = 90
MAX_TABLE_ROWS = 12
MAX_TABLE_COLS = 6
MAX_TABLE_HEADER_CELL_CHARS = 120
MAX_TABLE_BODY_CELL_CHARS = 180
EMU_PER_INCH = 914400
_UUID_RE = re.compile(
    r"\b[0-9a-f]{8}-[0-9a-f]{4}-[1-5][0-9a-f]{3}-[89ab][0-9a-f]{3}-[0-9a-f]{12}\b",
    re.IGNORECASE,
)


def _slugify(value: str) -> str:
    normalized = re.sub(r"[^A-Za-z0-9._-]+", "_", (value or "").strip())
    normalized = normalized.strip("._-")
    return normalized or "rag_slides"


def _truncate_footer_label(text: str, limit: int = 200) -> str:
    value = (text or "").strip()
    if len(value) <= limit:
        return value
    return value[: max(0, limit - len(_ELLIPSIS))].rstrip() + _ELLIPSIS


def _truncate_text(text: str, limit: int) -> str:
    value = (text or "").strip()
    if len(value) <= limit:
        return value
    return value[: max(0, limit - len(_ELLIPSIS))].rstrip() + _ELLIPSIS


def _read_uri_bytes(uri: str) -> bytes:
    if not uri.startswith("gs://"):
        return Path(uri).read_bytes()
    try:
        from google.cloud import storage

        _, _, rest = uri.partition("gs://")
        bucket_name, _, blob_name = rest.partition("/")
        if not bucket_name or not blob_name:
            raise ValueError(f"Invalid gs:// URI: {uri}")
        blob = storage.Client().bucket(bucket_name).blob(blob_name)
        return blob.download_as_bytes()
    except Exception:
        return subprocess.check_output(["gsutil", "cat", uri])


def _caption_is_unhelpful(caption: str) -> bool:
    text = (caption or "").strip()
    if not text:
        return True
    normalized = " ".join(text.lower().split())
    if _UUID_RE.search(normalized):
        return True
    if "from source" in normalized or "source page" in normalized:
        return True
    return bool(re.fullmatch(r"(table|tabella|figure|image|chart|asset)\s*[\w-]*", normalized))


def _table_spans_from_structured_obj(obj: Dict[str, Any]) -> list[dict[str, int]]:
    spans: list[dict[str, int]] = []
    explicit_spans = obj.get("spans")
    if isinstance(explicit_spans, list):
        for raw in explicit_spans:
            if not isinstance(raw, dict):
                continue
            try:
                row = int(raw.get("row", -1))
                col = int(raw.get("col", -1))
                row_span = int(raw.get("row_span", raw.get("rowSpan", 1)))
                col_span = int(raw.get("col_span", raw.get("colSpan", 1)))
            except Exception:
                continue
            if row < 0 or col < 0 or row_span < 1 or col_span < 1:
                continue
            if row_span > 1 or col_span > 1:
                spans.append({"row": row, "col": col, "row_span": row_span, "col_span": col_span})
        if spans:
            return spans

    rows_with_cells: list[dict[str, Any]] = []
    for row in obj.get("header_rows", []):
        if isinstance(row, dict):
            rows_with_cells.append(row)
    for row in obj.get("body_rows", []):
        if isinstance(row, dict):
            rows_with_cells.append(row)
    if not rows_with_cells:
        return []
    for r_idx, row in enumerate(rows_with_cells):
        cells = row.get("cells")
        if not isinstance(cells, list):
            continue
        for c_idx, cell in enumerate(cells):
            if not isinstance(cell, dict):
                continue
            try:
                row_span = int(cell.get("row_span", cell.get("rowSpan", 1)))
                col_span = int(cell.get("col_span", cell.get("colSpan", 1)))
            except Exception:
                continue
            if row_span > 1 or col_span > 1:
                spans.append({"row": r_idx, "col": c_idx, "row_span": row_span, "col_span": col_span})
    return spans


def _normalize_rows(rows: list[list[Any]]) -> list[list[str]]:
    out: list[list[str]] = []
    for row in rows:
        if not isinstance(row, list):
            continue
        normalized_row = [" ".join(str(c or "").replace("\n", " ").split()) for c in row]
        if any(cell for cell in normalized_row):
            out.append(normalized_row)
    return out


def _load_table_payload_from_visual(visual: Dict[str, Any]) -> dict[str, Any] | None:
    def _out(rows: list[list[Any]], table_obj: dict[str, Any] | None = None) -> dict[str, Any] | None:
        normalized_rows = _normalize_rows(rows)
        if not normalized_rows:
            return None
        spans = _table_spans_from_structured_obj(table_obj) if isinstance(table_obj, dict) else []
        return {"rows": normalized_rows, "spans": spans}

    inline_rows = visual.get("table_rows")
    if isinstance(inline_rows, list) and inline_rows:
        out_inline = _out(inline_rows)
        if out_inline:
            return out_inline

    table_obj = visual.get("table")
    if isinstance(table_obj, dict) and isinstance(table_obj.get("rows"), list):
        out_table = _out(table_obj.get("rows", []), table_obj)
        if out_table:
            return out_table

    files = visual.get("files") if isinstance(visual.get("files"), dict) else {}
    local_table_path = str(files.get("table_local_path") or "").strip()
    if local_table_path:
        try:
            obj = json.loads(Path(local_table_path).read_text(encoding="utf-8", errors="replace"))
            rows = obj.get("rows") if isinstance(obj, dict) else None
        except Exception:
            rows = None
            obj = None
        if isinstance(rows, list) and rows:
            out_local = _out(rows, obj if isinstance(obj, dict) else None)
            if out_local:
                return out_local

    table_uri = str(files.get("table_uri") or "").strip()
    if not table_uri:
        return None
    try:
        obj = json.loads(_read_uri_bytes(table_uri).decode("utf-8", errors="replace"))
    except Exception:
        return None
    rows = obj.get("rows") if isinstance(obj, dict) else None
    if not isinstance(rows, list) or not rows:
        return None
    return _out(rows, obj if isinstance(obj, dict) else None)


def _drop_empty_columns(rows: list[list[str]]) -> tuple[list[list[str]], list[int]]:
    if not rows:
        return rows, []
    max_cols = max((len(r) for r in rows), default=0)
    if max_cols <= 0:
        return rows, []
    padded = [list(r[:max_cols]) + [""] * (max_cols - len(r[:max_cols])) for r in rows]
    keep_cols = []
    for c in range(max_cols):
        if any((row[c] or "").strip() for row in padded):
            keep_cols.append(c)
    if not keep_cols:
        return [], []
    return [[row[c] for c in keep_cols] for row in padded], keep_cols


def _fit_table_rows(
    rows: list[list[str]], spans: list[dict[str, int]] | None = None
) -> tuple[list[list[str]], list[dict[str, int]]]:
    if not rows:
        return rows, []
    non_empty_rows = [row for row in rows if any((c or "").strip() for c in row)]
    if not non_empty_rows:
        return [], []
    trimmed_rows, keep_cols = _drop_empty_columns(non_empty_rows)
    if not trimmed_rows:
        return [], []
    source_max_cols = max((len(r) for r in trimmed_rows), default=1)
    max_cols = max(1, min(MAX_TABLE_COLS, source_max_cols))
    keep_cols = keep_cols[:max_cols]
    out_rows: list[list[str]] = []
    for i, row in enumerate(trimmed_rows[:MAX_TABLE_ROWS]):
        clipped = list(row[:max_cols]) + [""] * max(0, max_cols - len(row[:max_cols]))
        if i == 0 and clipped:
            prev_header = ""
            for j, cell in enumerate(clipped):
                value = cell.strip()
                if value:
                    prev_header = value
                    continue
                if prev_header:
                    clipped[j] = prev_header
                else:
                    clipped[j] = "—"
        cell_limit = MAX_TABLE_HEADER_CELL_CHARS if i == 0 else MAX_TABLE_BODY_CELL_CHARS
        out_rows.append([_truncate_text(cell, cell_limit) for cell in clipped])
    if not spans:
        return out_rows, []

    out_spans: list[dict[str, int]] = []
    max_row = len(out_rows)
    max_col = len(out_rows[0]) if out_rows else 0
    for span in spans:
        try:
            row = int(span.get("row", -1))
            col = int(span.get("col", -1))
            row_span = max(1, int(span.get("row_span", 1)))
            col_span = max(1, int(span.get("col_span", 1)))
        except Exception:
            continue
        if row < 0 or col < 0:
            continue
        if row >= MAX_TABLE_ROWS or row >= max_row or col not in keep_cols:
            continue
        new_col = keep_cols.index(col)
        row_end = min(max_row - 1, row + row_span - 1)
        col_end = min(max_col - 1, new_col + col_span - 1)
        new_row_span = row_end - row + 1
        new_col_span = col_end - new_col + 1
        if new_row_span > 1 or new_col_span > 1:
            out_spans.append({"row": row, "col": new_col, "row_span": new_row_span, "col_span": new_col_span})
    return out_rows, out_spans


def _visual_locator_text(visual: Dict[str, Any]) -> str:
    locator = visual.get("locator") if isinstance(visual.get("locator"), dict) else {}
    doc_id = str(visual.get("doc_id") or "").strip()
    page = locator.get("page_start", visual.get("page"))
    try:
        page_no = int(page)
    except Exception:
        page_no = 0
    table_idx_raw = locator.get("table_index", locator.get("table_idx", visual.get("table_index")))
    try:
        table_idx = int(table_idx_raw)
    except Exception:
        table_idx = None
    parts: list[str] = []
    if doc_id:
        parts.append(doc_id)
    if page_no > 0:
        parts.append(f"p.{page_no}")
    if table_idx is not None and table_idx >= 0:
        parts.append(f"t#{table_idx}")
    if not parts:
        parts.append("source")
    return ", ".join(parts)


def _table_short_title(rows: list[list[str]] | None) -> str:
    if not rows:
        return ""
    header = rows[0] if rows else []
    parts = [c.strip() for c in header[:3] if str(c).strip() and str(c).strip() != "—"]
    return _truncate_text(" | ".join(parts), 70) if parts else ""


def _safe_visual_caption(visual: Dict[str, Any]) -> str:
    caption = str(visual.get("caption") or "").strip()
    if caption and not _caption_is_unhelpful(caption):
        return caption
    asset_type = str(visual.get("asset_type") or "asset").strip().lower() or "asset"
    locator_text = _visual_locator_text(visual)
    if asset_type == "table":
        rows = visual.get("_table_rows")
        title = _table_short_title(rows if isinstance(rows, list) else None)
        if title:
            return f"Table ({locator_text}): {title}"
        return f"Table ({locator_text})"
    return f"{asset_type.title()} ({locator_text})"


def _table_text_overflow_probable(rows: list[list[str]], width, height, body_font_size: int) -> bool:
    if not rows:
        return False
    max_rows = len(rows)
    max_cols = max((len(r) for r in rows), default=1)
    col_w_in = (float(width) / EMU_PER_INCH) / max(1, max_cols)
    row_h_in = (float(height) / EMU_PER_INCH) / max(1, max_rows)
    char_per_line = max(6, int(col_w_in * 10.0 * (9.0 / max(1.0, float(body_font_size)))))
    max_lines = max(1, int((row_h_in * 72.0) / (max(7.0, float(body_font_size)) * 1.45)))
    over = 0
    total = 0
    for r, row in enumerate(rows):
        limit = char_per_line * max_lines
        if r == 0:
            limit = max(limit, int(limit * 1.2))
        for cell in row:
            text = (cell or "").strip()
            if not text:
                continue
            total += 1
            if len(text) > limit:
                over += 1
    if total <= 0:
        return False
    return (over / total) >= 0.28


def _add_table_visual(
    slide,
    rows: list[list[str]],
    spans: list[dict[str, int]] | None,
    left,
    top,
    width,
    height,
) -> tuple[bool, str]:
    normalized, normalized_spans = _fit_table_rows(rows, spans)
    if not normalized:
        return False, "empty_table"
    height_in = max(0.5, float(height) / EMU_PER_INCH)
    max_rows_by_height = max(2, int(height_in / 0.24))
    normalized = normalized[:max_rows_by_height]
    normalized_spans = [
        s
        for s in normalized_spans
        if s["row"] < len(normalized) and s["col"] < (len(normalized[0]) if normalized else 0)
    ]
    normalized_spans = [
        {
            "row": s["row"],
            "col": s["col"],
            "row_span": min(s["row_span"], len(normalized) - s["row"]),
            "col_span": min(s["col_span"], len(normalized[0]) - s["col"]),
        }
        for s in normalized_spans
    ]
    normalized_spans = [s for s in normalized_spans if s["row_span"] > 1 or s["col_span"] > 1]

    max_rows = len(normalized)
    max_cols = max(len(r) for r in normalized)
    table_shape = slide.shapes.add_table(max_rows, max_cols, left, top, width, height)
    table = table_shape.table
    col_w = int(float(width) / max_cols)
    for c in range(max_cols):
        table.columns[c].width = col_w
    row_h = int(float(height) / max_rows)
    for r in range(max_rows):
        table.rows[r].height = row_h
    if max_rows >= 10:
        body_font_size = 8
    elif max_rows >= 8:
        body_font_size = 8
    else:
        body_font_size = 9
    if _table_text_overflow_probable(normalized, width, height, body_font_size):
        body_font_size = 7

    for r in range(max_rows):
        for c in range(max_cols):
            cell = table.cell(r, c)
            cell.text = normalized[r][c]
            tf = cell.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            cell.margin_left = Pt(3)
            cell.margin_right = Pt(3)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = FONT_NAME
                    run.font.size = Pt(10 if r == 0 else body_font_size)
                    run.font.bold = bool(r == 0 or c == 0)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
    if not _KEYNOTE_SAFE:
        for span in normalized_spans:
            try:
                row = span["row"]
                col = span["col"]
                row2 = row + span["row_span"] - 1
                col2 = col + span["col_span"] - 1
                if row2 > row or col2 > col:
                    table.cell(row, col).merge(table.cell(row2, col2))
            except Exception:
                continue
    if _table_text_overflow_probable(normalized, width, height, body_font_size):
        return True, "overflow_risk"
    return True, "ok"


def _add_picture_contain(slide, image_path: str, left, top, width, height) -> None:
    """
    Place image in a slot preserving aspect ratio (contain + center), avoiding stretch.
    """
    try:
        with Image.open(image_path) as im:
            img_w, img_h = im.size
    except Exception:
        slide.shapes.add_picture(image_path, left, top, width, height)
        return

    if img_w <= 0 or img_h <= 0:
        slide.shapes.add_picture(image_path, left, top, width, height)
        return

    slot_w = float(width)
    slot_h = float(height)
    img_ratio = float(img_w) / float(img_h)
    slot_ratio = slot_w / slot_h if slot_h > 0 else img_ratio

    if img_ratio >= slot_ratio:
        fit_w = slot_w
        fit_h = fit_w / img_ratio
    else:
        fit_h = slot_h
        fit_w = fit_h * img_ratio

    fit_left = int(float(left) + (slot_w - fit_w) / 2.0)
    fit_top = int(float(top) + (slot_h - fit_h) / 2.0)
    slide.shapes.add_picture(image_path, fit_left, fit_top, int(fit_w), int(fit_h))


def build_pptx_from_slide_plan(
    slide_plan: Dict[str, Any],
    *,
    out_dir: str = "out",
    filename: Optional[str] = None,
) -> Path:
    slides = slide_plan.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("slide_plan['slides'] must be a non-empty list")

    raw_title = str(slide_plan.get("title") or "").strip()
    if raw_title:
        deck_title = raw_title
    else:
        outline_used = slide_plan.get("outline_used")
        if isinstance(outline_used, dict):
            outline_title = str(outline_used.get("title") or "").strip()
            deck_title = outline_title or "RAG Slides"
        else:
            deck_title = "RAG Slides"

    if filename is not None:
        generated_name = filename
        if not generated_name.lower().endswith(".pptx"):
            generated_name += ".pptx"
    else:
        payload = json.dumps(slide_plan, sort_keys=True, ensure_ascii=False, separators=(",", ":"))
        sha = hashlib.sha1(payload.encode("utf-8")).hexdigest()[:10]
        generated_name = f"{_slugify(deck_title)}_{sha}.pptx"

    out_path = Path(out_dir) / generated_name
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    # Force widescreen 16:9.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, raw_slide in enumerate(slides, start=1):
        slide_data = raw_slide if isinstance(raw_slide, dict) else {}
        title = _truncate_text(str(slide_data.get("title") or f"Slide {idx}"), MAX_TITLE_CHARS)
        bullets_raw = slide_data.get("bullets")
        bullets = (
            [_truncate_text(str(b), MAX_BULLET_CHARS) for b in bullets_raw if str(b).strip()]
            if isinstance(bullets_raw, list)
            else []
        )
        bullets = bullets[:MAX_BULLETS_PER_SLIDE]
        citations_raw = slide_data.get("citations")
        citations = [str(c).strip() for c in citations_raw if str(c).strip()] if isinstance(citations_raw, list) else []
        sources_raw = slide_data.get("sources")
        sources = [s for s in sources_raw if isinstance(s, dict)] if isinstance(sources_raw, list) else []
        visuals_raw = slide_data.get("visuals")
        visuals = [v for v in visuals_raw if isinstance(v, dict)] if isinstance(visuals_raw, list) else []

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.left = Inches(0.6)
        title_shape.top = Inches(0.3)
        title_shape.width = Inches(12.1)
        title_shape.height = Inches(0.9)
        for p in title_shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = FONT_NAME
                run.font.size = Pt(TITLE_FONT_SIZE_PT)

        body_shape = slide.placeholders[1]
        body_shape.left = Inches(0.7)
        body_shape.top = Inches(1.25)
        body_shape.width = Inches(11.9)
        if visuals:
            body_shape.height = Inches(1.55)
            bullets = bullets[:3]
        else:
            body_shape.height = Inches(4.8)
        text_frame = body_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        for bullet_idx, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if bullet_idx == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.space_after = Pt(6)
            for run in paragraph.runs:
                run.font.name = FONT_NAME
                run.font.size = Pt(BODY_FONT_SIZE_PT)

        visual_notes: list[str] = []
        if visuals:
            valid_visuals = []
            for v_i, v in enumerate(visuals[:2], start=1):
                rp = str(v.get("render_path") or "").strip()
                table_payload = None
                if str(v.get("asset_type") or "").strip().lower() == "table":
                    table_payload = _load_table_payload_from_visual(v)
                if table_payload is not None:
                    v = dict(v)
                    v["_table_rows"] = table_payload.get("rows")
                    v["_table_spans"] = table_payload.get("spans")
                    valid_visuals.append(v)
                elif rp and Path(rp).exists():
                    valid_visuals.append(v)
                else:
                    visual_notes.append(f"MISSING_ASSET asset_id={v.get('asset_id')}")
                locator_text = _visual_locator_text(v)
                v_type = str(v.get("asset_type") or "").strip().lower()
                if not _KEYNOTE_SAFE:
                    if v_type == "table":
                        visual_notes.append(
                            f"TABLE_{v_i} asset_id={v.get('asset_id')} doc_id={v.get('doc_id')} locator={locator_text}"
                        )
                    else:
                        visual_notes.append(
                            f"FIGURE_{v_i} asset_id={v.get('asset_id')} doc_id={v.get('doc_id')} locator={locator_text}"
                        )

            # If slide carries no textual content and visuals are all missing, fail deterministically.
            if not bullets and not valid_visuals:
                raise ValueError(f"insufficient assets on slide {idx}: all visuals missing")

            if len(valid_visuals) == 1:
                v = valid_visuals[0]
                img_left, img_top, img_w, img_h = Inches(0.9), Inches(3.15), Inches(11.5), Inches(2.15)
                rows = v.get("_table_rows")
                if isinstance(rows, list):
                    if _KEYNOTE_SAFE and not _KEYNOTE_KEEP_TABLES and str(v.get("render_path") or "").strip() and Path(str(v["render_path"])).exists():
                        _add_picture_contain(slide, str(v["render_path"]), img_left, img_top, img_w, img_h)
                    else:
                        spans = v.get("_table_spans") if isinstance(v.get("_table_spans"), list) else []
                        rendered, reason = _add_table_visual(slide, rows, spans, img_left, img_top, img_w, img_h)
                        if not rendered and str(v.get("render_path") or "").strip() and Path(str(v["render_path"])).exists():
                            _add_picture_contain(slide, str(v["render_path"]), img_left, img_top, img_w, img_h)
                            if not _KEYNOTE_SAFE:
                                visual_notes.append(f"TABLE_RENDER_FALLBACK asset_id={v.get('asset_id')} reason={reason}")
                else:
                    _add_picture_contain(slide, str(v["render_path"]), img_left, img_top, img_w, img_h)
                cap = _safe_visual_caption(v)
                cap_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.35))
                cap_tf = cap_box.text_frame
                cap_tf.clear()
                cp = cap_tf.paragraphs[0]
                cp.text = _truncate_text(cap, 180)
                for run in cp.runs:
                    run.font.name = FONT_NAME
                    run.font.size = Pt(10)
            elif len(valid_visuals) >= 2:
                slots = [
                    (Inches(0.8), Inches(3.05), Inches(5.8), Inches(2.0)),
                    (Inches(6.75), Inches(3.05), Inches(5.8), Inches(2.0)),
                ]
                for v, slot in zip(valid_visuals[:2], slots):
                    rows = v.get("_table_rows")
                    if isinstance(rows, list):
                        if _KEYNOTE_SAFE and not _KEYNOTE_KEEP_TABLES and str(v.get("render_path") or "").strip() and Path(str(v["render_path"])).exists():
                            _add_picture_contain(slide, str(v["render_path"]), *slot)
                        else:
                            spans = v.get("_table_spans") if isinstance(v.get("_table_spans"), list) else []
                            rendered, reason = _add_table_visual(slide, rows, spans, *slot)
                            if not rendered and str(v.get("render_path") or "").strip() and Path(str(v["render_path"])).exists():
                                _add_picture_contain(slide, str(v["render_path"]), *slot)
                                if not _KEYNOTE_SAFE:
                                    visual_notes.append(f"TABLE_RENDER_FALLBACK asset_id={v.get('asset_id')} reason={reason}")
                    else:
                        _add_picture_contain(slide, str(v["render_path"]), *slot)
                    cap = _safe_visual_caption(v)
                    cap_box = slide.shapes.add_textbox(slot[0], slot[1] + slot[3] + Inches(0.05), slot[2], Inches(0.32))
                    cap_tf = cap_box.text_frame
                    cap_tf.clear()
                    cp = cap_tf.paragraphs[0]
                    cp.text = _truncate_text(cap, 90)
                    for run in cp.runs:
                        run.font.name = FONT_NAME
                        run.font.size = Pt(9)

        footer_lines: list[str] = []
        if citations:
            src_by_ref = {}
            for src in sources:
                ref_id = str(src.get("ref_id") or "").strip()
                if ref_id and ref_id not in src_by_ref:
                    src_by_ref[ref_id] = src

            grouped: Dict[str, Dict[str, Any]] = {}
            group_order: list[str] = []
            for c in citations:
                src = src_by_ref.get(c)
                if not src:
                    continue
                item_id = str(src.get("item_id") or "").strip()
                doc_id = str(src.get("doc_id") or "").strip()
                doc_citation = str(src.get("doc_citation") or "").strip()
                doc_filename = str(src.get("doc_filename") or "").strip()
                label = _truncate_footer_label(doc_citation or doc_filename or item_id or c)
                group_key = doc_id or f"fallback::{label}"
                if group_key not in grouped:
                    grouped[group_key] = {"label": label, "tokens": []}
                    group_order.append(group_key)
                if c not in grouped[group_key]["tokens"]:
                    grouped[group_key]["tokens"].append(c)

            footer_lines = [
                f"{grouped[key]['label']} ({', '.join(grouped[key]['tokens'])})"
                for key in group_order
                if grouped[key]["tokens"]
            ]

        if footer_lines:
            slide_width_in = float(prs.slide_width) / 914400.0
            slide_height_in = float(prs.slide_height) / 914400.0
            footer_height = min(1.2, max(0.3, 0.16 * len(footer_lines)))
            footer_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(max(0.1, slide_height_in - footer_height - 0.1)),
                Inches(max(1.0, slide_width_in - 1.0)),
                Inches(footer_height),
            )
            footer_frame = footer_box.text_frame
            footer_frame.clear()
            footer_frame.word_wrap = True
            footer_frame.auto_size = MSO_AUTO_SIZE.NONE if _KEYNOTE_SAFE else MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            footer_p = footer_frame.paragraphs[0]
            footer_p.text = "\n".join(footer_lines)
            footer_p.alignment = PP_ALIGN.LEFT
            for run in footer_p.runs:
                run.font.name = FONT_NAME
                run.font.size = Pt(FOOTER_FONT_SIZE_PT)

        notes_parts = []
        speaker_notes = slide_data.get("speaker_notes")
        if speaker_notes is not None:
            notes_text = str(speaker_notes).strip()
            if notes_text:
                notes_parts.append(notes_text)

        if citations:
            src_by_ref = {}
            for src in sources:
                ref_id = str(src.get("ref_id") or "").strip()
                if ref_id and ref_id not in src_by_ref:
                    src_by_ref[ref_id] = src

            source_lines = ["Sources:"]
            for c in citations:
                src = src_by_ref.get(c)
                if not src:
                    source_lines.append(f"{c} (missing source)")
                    continue
                item_id = str(src.get("item_id") or "").strip()
                doc_citation = str(src.get("doc_citation") or "").strip()
                doc_filename = str(src.get("doc_filename") or "").strip()
                page_start = int(src.get("page_start", 0))
                page_end = int(src.get("page_end", 0))
                if page_start == page_end:
                    page_text = f"p. {page_start}"
                else:
                    page_text = f"p. {page_start}-{page_end}"
                label = doc_citation or doc_filename or item_id or c
                source_lines.append(f"{label} ({c} {page_text})")
            notes_parts.append("\n".join(source_lines))

        if not _KEYNOTE_SAFE:
            if notes_parts:
                slide.notes_slide.notes_text_frame.text = "\n\n".join(notes_parts)
                if visual_notes:
                    slide.notes_slide.notes_text_frame.text += "\n\n" + "\n".join(visual_notes)
            elif visual_notes:
                slide.notes_slide.notes_text_frame.text = "\n".join(visual_notes)

    prs.save(str(out_path))
    return out_path
