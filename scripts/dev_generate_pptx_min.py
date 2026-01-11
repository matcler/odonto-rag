import os, requests
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches
from qdrant_client import QdrantClient

PROJECT_ID = os.environ["PROJECT_ID"]
LOCATION = os.environ["LOCATION"]
QDRANT_URL = os.environ["QDRANT_URL"]
QDRANT_API_KEY = os.environ["QDRANT_API_KEY"]
COL = os.environ["QDRANT_COLLECTION"]

def embed(text):
    token = os.popen("gcloud auth print-access-token").read().strip()
    url = f"https://{LOCATION}-aiplatform.googleapis.com/v1/projects/{PROJECT_ID}/locations/{LOCATION}/publishers/google/models/gemini-embedding-001:predict"
    r = requests.post(url,
        headers={"Authorization": f"Bearer {token}", "Content-Type":"application/json"},
        json={"instances":[{"content": text}]},
        timeout=60
    )
    r.raise_for_status()
    return r.json()["predictions"][0]["embeddings"]["values"]

def retrieve(query, limit=25):
    q = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY)
    vec = embed(query)
    return q.query_points(collection_name=COL, query=vec, limit=limit).points

def make_pptx(title, query, out_path="odonto_output.pptx"):
    hits = retrieve(query, limit=30)

    # group by source file for simple “citation density”
    by_source = defaultdict(list)
    for h in hits:
        p = h.payload or {}
        by_source[p.get("source_object_path","unknown")].append(h)

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Auto-generated from your indexed documents (Qdrant + embeddings)"

    # A simple fixed outline (you can change)
    outline = [
        "Background & definitions",
        "Classification of adhesive systems",
        "Etch-and-rinse systems",
        "Self-etch systems",
        "Universal / multi-mode adhesives",
        "Clinical performance outcomes",
        "Technique sensitivity & common errors",
        "Substrate considerations (enamel vs dentin)",
        "Longevity, failure modes, marginal integrity",
        "Practical recommendations",
        "Take-home messages",
        "References (sources)"
    ]

    # helper to add a bullet slide from top hits
    def add_bullet_slide(slide_title, selected_hits):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = slide_title
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()

        # bullets: take first 3 chunks, truncate
        for h in selected_hits[:3]:
            p = h.payload or {}
            txt = (p.get("text") or "").replace("\n", " ").strip()
            if not txt:
                continue
            bullet = txt[:180].rstrip() + ("…" if len(txt) > 180 else "")
            para = tf.add_paragraph() if tf.text else tf.paragraphs[0]
            para.text = bullet

        # speaker notes with citations
        notes = s.notes_slide.notes_text_frame
        notes.text = "Sources used:\n"
        for h in selected_hits[:5]:
            p = h.payload or {}
            notes.text += f"- {p.get('source_object_path')} (pp. {p.get('page_start')}-{p.get('page_end')}) score={h.score:.3f}\n"

    # naive selection: use best hits for each slide (works fine for a first version)
    # Tomorrow we can do smarter mapping slide->query
    for t in outline[:-1]:
        add_bullet_slide(t, hits)

    # References slide
    ref = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    ref.shapes.title.text = "References (indexed sources)"
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8.5)
    height = Inches(4.5)
    box = ref.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = ""
    for src in sorted(by_source.keys()):
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = src

    prs.save(out_path)
    print("✅ Saved:", out_path)

if __name__ == "__main__":
    # change these two lines for your request
    TITLE = "Dental Adhesive Systems — Overview & Clinical Evidence"
    QUERY = "dental adhesive systems etch-and-rinse self-etch universal clinical performance meta-analysis"
    make_pptx(TITLE, QUERY, out_path="odonto_adhesives.pptx")


def main():
    # Minimal wrapper to keep CLI behavior stable during refactor
    # The script logic should already run under __main__ below.
    pass

if __name__ == "__main__":
    main()
